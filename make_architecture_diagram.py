"""
Generate SpiceFusionNet architecture diagram and embed in a PPTX slide.
Output: outputs/architecture_diagram.pptx
"""
import sys
sys.path.insert(0, "D:/SpiceNet")

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

# ── Color Palette ─────────────────────────────────────────────────────────────
C_INPUT  = "#1B2631"   # near-black  — input/output
C_CNN    = "#1A6B8A"   # deep teal   — CNN branch
C_TEX    = "#784212"   # earthy brown— texture branch
C_COL    = "#1E8449"   # forest green— color branch
C_FUSE   = "#6C3483"   # purple      — fusion
C_HEAD   = "#922B21"   # deep red    — heads / loss
C_BG     = "#F0F3F4"   # light grey  background
C_SECT   = "#ECF0F1"   # section bg

TEXT_W   = "#FFFFFF"
TEXT_D   = "#1A1A2E"

FIG_W, FIG_H = 24, 15.5

def box(ax, cx, cy, w, h, fc, lines, fs=9, tc=TEXT_W,
        rad=0.28, bold0=True, zorder=4, lw=1.8):
    p = FancyBboxPatch((cx-w/2, cy-h/2), w, h,
                       boxstyle=f"round,pad=0,rounding_size={rad}",
                       facecolor=fc, edgecolor="white", linewidth=lw,
                       zorder=zorder)
    ax.add_patch(p)
    n = len(lines)
    step = h / (n + 1)
    for i, txt in enumerate(lines):
        fw = "bold" if (i == 0 and bold0) else "normal"
        ax.text(cx, cy + h/2 - step*(i+1), txt, ha="center", va="center",
                fontsize=fs, color=tc, fontweight=fw, zorder=zorder+1,
                linespacing=1.3)

def badge(ax, cx, cy, txt, fc, fs=7.5, zorder=6):
    ax.text(cx, cy, txt, ha="center", va="center", fontsize=fs,
            color="white", fontweight="bold", zorder=zorder,
            bbox=dict(boxstyle="round,pad=0.28", facecolor=fc,
                      edgecolor="white", linewidth=1.1, alpha=0.93))

def arr(ax, x0, y0, x1, y1, fc=C_INPUT, lw=1.9, ms=15, rad=0.0):
    ax.annotate("", xy=(x1, y1), xytext=(x0, y0),
                arrowprops=dict(arrowstyle="-|>", color=fc, lw=lw,
                                mutation_scale=ms,
                                connectionstyle=f"arc3,rad={rad}"),
                zorder=3)

def sect_bg(ax, cx, cy, w, h, fc, label="", lfs=8.5):
    p = FancyBboxPatch((cx-w/2, cy-h/2), w, h,
                       boxstyle="round,pad=0,rounding_size=0.6",
                       facecolor=fc, edgecolor=fc,
                       linewidth=2.0, alpha=0.10, zorder=1, linestyle="--")
    ax.add_patch(p)
    if label:
        ax.text(cx - w/2 + 0.2, cy + h/2 - 0.18, label,
                fontsize=lfs, color=fc, alpha=0.75, fontweight="bold",
                va="top", ha="left", zorder=2)

def phase_tag(ax, cx, cy, txt, fc):
    ax.text(cx, cy, txt, ha="center", va="center", fontsize=8,
            color="white", fontweight="bold", zorder=8,
            bbox=dict(boxstyle="round,pad=0.35", facecolor=fc,
                      edgecolor="white", linewidth=1.3, alpha=0.95))

# ── Layout coordinates ────────────────────────────────────────────────────────
#   X:  CNN=5.0  TEX=14.0  COL=20.0  FUSE=11.0
#   Y (top=high): INPUT=14.6  FEAT=12.8  SPLIT=10.5  BRANCH=8.2  FUSE=5.7  OUT=3.4  FINAL=1.6

X_CNN  =  5.0;  X_TEX  = 14.0;  X_COL  = 20.0;  X_FUSE = 11.5
Y_IN   = 14.6;  Y_FE   = 12.8;  Y_SP   = 10.5;  Y_BR   =  8.2
Y_FU   =  5.7;  Y_OUT  =  3.4;  Y_FIN  =  1.6

fig, ax = plt.subplots(figsize=(FIG_W, FIG_H))
ax.set_xlim(0, FIG_W);  ax.set_ylim(0, FIG_H)
ax.set_facecolor(C_BG);  fig.patch.set_facecolor(C_BG);  ax.axis("off")

# ── Title ─────────────────────────────────────────────────────────────────────
ax.text(FIG_W/2, 15.15, "SpiceFusionNet  —  Multi-Modal Architecture",
        ha="center", va="center", fontsize=19, fontweight="bold", color=C_INPUT)
ax.text(FIG_W/2, 14.72,
        "EfficientNet-B4   x   Texture (LBP+GLCM)   x   Color (HSV)   -->   AttentionFusion   -->   11-class",
        ha="center", va="center", fontsize=10, color="#566573")

# ── Section backgrounds ───────────────────────────────────────────────────────
sect_bg(ax,  5.0, 10.85, 8.0,  9.0, C_CNN,  "CNN Branch")
sect_bg(ax, 14.0, 10.85, 6.5,  9.0, C_TEX,  "Texture Branch")
sect_bg(ax, 20.0, 10.85, 6.0,  9.0, C_COL,  "Color Branch")
sect_bg(ax, 11.5,  3.75, 23.0, 5.0, C_FUSE, "Fusion & Output (Phase 3)")

# ── INPUT IMAGE ───────────────────────────────────────────────────────────────
box(ax, X_FUSE, Y_IN, 6.0, 1.0, C_INPUT,
    ["Input Image", "224 x 224 x 3"], fs=11)
# note for hand-crafted features
ax.text(18.5, 13.5, "at 512x512  (before augmentation)",
        fontsize=8, color=C_TEX, fontstyle="italic", ha="center")

# Input splits -> CNN, TEX, COL
arr(ax, 9.5, Y_IN-0.5,  X_CNN+1.5, Y_FE+0.6,  fc=C_CNN)
arr(ax, 11.5, Y_IN-0.5, X_TEX,     Y_FE+0.6,  fc=C_TEX)
arr(ax, 13.5, Y_IN-0.5, X_COL-1.0, Y_FE+0.55, fc=C_COL)

# ── EfficientNet-B4 ───────────────────────────────────────────────────────────
box(ax, X_CNN, Y_FE, 4.2, 1.5, C_CNN,
    ["EfficientNet-B4", "(ImageNet pretrained)", "Global Average Pool"], fs=9.5)
badge(ax, X_CNN, Y_FE-0.9, "-> 1792-d", C_CNN)
arr(ax, X_CNN, Y_FE-1.15, X_CNN, Y_SP+0.5, fc=C_CNN)

# ── LBP ──────────────────────────────────────────────────────────────────────
box(ax, X_TEX-1.3, Y_FE, 2.5, 1.4, C_TEX,
    ["LBP", "P=8, R=1", "uniform"], fs=9)
badge(ax, X_TEX-1.3, Y_FE-0.85, "10-d", C_TEX)

# ── GLCM ─────────────────────────────────────────────────────────────────────
box(ax, X_TEX+1.3, Y_FE, 2.5, 1.4, C_TEX,
    ["GLCM", "6 props x 2d x 4 ang"], fs=9)
badge(ax, X_TEX+1.3, Y_FE-0.85, "48-d", C_TEX)

# LBP+GLCM -> concat
arr(ax, X_TEX-1.3, Y_FE-1.1, X_TEX, Y_SP+0.45, fc=C_TEX)
arr(ax, X_TEX+1.3, Y_FE-1.1, X_TEX, Y_SP+0.45, fc=C_TEX)

box(ax, X_TEX, Y_SP, 3.2, 0.9, C_TEX,
    ["concat: LBP + GLCM", "58-d texture vector"], fs=8.5)
arr(ax, X_TEX, Y_SP-0.45, X_TEX, Y_BR+0.55, fc=C_TEX)

# ── HSV Histogram ─────────────────────────────────────────────────────────────
box(ax, X_COL, Y_FE, 3.8, 1.5, C_COL,
    ["HSV Histogram", "H: 36 bins", "S: 32 bins,  V: 32 bins"], fs=9)
badge(ax, X_COL, Y_FE-0.9, "100-d", C_COL)
arr(ax, X_COL, Y_FE-1.15, X_COL, Y_BR+0.55, fc=C_COL)

# ── CNN split: img_head + proj_head ──────────────────────────────────────────
# img_head (Phase 1)
box(ax, X_CNN-2.0, Y_SP, 3.0, 1.0, C_HEAD,
    ["img_head  (Phase 1)", "Linear(1792->512->11)"], fs=8.5)
badge(ax, X_CNN-2.0, Y_SP-0.65, "11 logits  |  CE Loss", C_HEAD)
arr(ax, X_CNN-0.8, Y_FE-0.7, X_CNN-2.0, Y_SP+0.5, fc=C_HEAD, rad=-0.25)

# proj_head (Phase 2)
box(ax, X_CNN+2.2, Y_SP, 3.2, 1.0, C_CNN,
    ["proj_head  (Phase 2)", "Linear(1792->512->128)"], fs=8.5)
badge(ax, X_CNN+2.2, Y_SP-0.65, "128-d  L2-norm  |  SupCon", C_CNN)
arr(ax, X_CNN+0.8, Y_FE-0.7, X_CNN+2.2, Y_SP+0.5, fc=C_CNN, rad=0.25)

# Phase tags
phase_tag(ax, X_CNN-4.0, Y_SP+0.0, "Phase 1", C_HEAD)
phase_tag(ax, X_CNN-4.0, Y_SP-0.7, "Phase 2", C_CNN)
arr(ax, X_CNN-3.55, Y_SP+0.0, X_CNN-2.0-1.5, Y_SP+0.05, fc=C_HEAD, lw=1.2)
arr(ax, X_CNN-3.55, Y_SP-0.7, X_CNN+2.2-1.6, Y_SP-0.3,  fc=C_CNN, lw=1.2)

# CNN 1792 features continue down
arr(ax, X_CNN, Y_SP-0.15, X_CNN, Y_BR+0.55, fc=C_CNN)

# ── Branch MLPs ───────────────────────────────────────────────────────────────
box(ax, X_CNN, Y_BR, 3.6, 1.1, C_CNN,
    ["CNN Features", "1792-d"], fs=10)

box(ax, X_TEX, Y_BR, 3.4, 1.3, C_TEX,
    ["texture_branch", "MLP(58->128->256)", "BN -> ReLU x 2"], fs=9)
badge(ax, X_TEX, Y_BR-0.82, "256-d output", C_TEX)

box(ax, X_COL, Y_BR, 3.6, 1.3, C_COL,
    ["color_branch", "MLP(100->64->128)", "BN -> ReLU x 2"], fs=9)
badge(ax, X_COL, Y_BR-0.82, "128-d output", C_COL)

# Arrows down from branches
arr(ax, X_CNN, Y_BR-0.98, X_CNN,  Y_FU+0.88, fc=C_CNN)
arr(ax, X_TEX, Y_BR-1.07, X_TEX,  Y_FU+0.88, fc=C_TEX)
arr(ax, X_COL, Y_BR-1.07, X_COL,  Y_FU+0.88, fc=C_COL)

# Dim labels on arrows
badge(ax, X_CNN,  Y_BR-1.38, "1792-d", C_CNN, fs=7.2)
badge(ax, X_TEX,  Y_BR-1.42, "256-d",  C_TEX, fs=7.2)
badge(ax, X_COL,  Y_BR-1.42, "128-d",  C_COL, fs=7.2)

# ── Concat label ──────────────────────────────────────────────────────────────
ax.text(X_FUSE, Y_FU+1.02,
        "concat:  1792 + 256 + 128  =  2176-d",
        ha="center", va="center", fontsize=9.5, fontweight="bold", color=C_FUSE,
        bbox=dict(boxstyle="round,pad=0.35", facecolor=C_FUSE,
                  alpha=0.13, edgecolor=C_FUSE, linewidth=1.2))

# ── AttentionFusion ───────────────────────────────────────────────────────────
box(ax, X_FUSE, Y_FU, 8.5, 1.5, C_FUSE,
    ["AttentionFusion",
     "gate: Linear(2176->3) -> Softmax",
     "a_img*f_cnn  +  a_tex*f_tex  +  a_col*f_col"], fs=9.5)
badge(ax, X_FUSE, Y_FU-0.93, "->  2176-d  (attention-weighted)", C_FUSE)
arr(ax, X_FUSE, Y_FU-1.18, X_FUSE, Y_OUT+0.7, fc=C_FUSE)

# ── Fusion Head ───────────────────────────────────────────────────────────────
box(ax, X_FUSE, Y_OUT, 9.0, 1.3, C_HEAD,
    ["fusion_head",
     "Linear(2176->512) -> BN -> ReLU -> Dropout(0.4) -> Linear(512->11)",
     "Loss:  0.5 x CE  +  0.5 x SupCon"], fs=9.5)
badge(ax, X_FUSE, Y_OUT-0.85, "11 class logits", C_HEAD)
arr(ax, X_FUSE, Y_OUT-1.08, X_FUSE, Y_FIN+0.45, fc=C_HEAD)

# ── Final Output ──────────────────────────────────────────────────────────────
box(ax, X_FUSE, Y_FIN, 8.5, 0.85, C_INPUT,
    ["Predicted Spice Class     Top-1: 99.00%   |   Top-5: 99.95%"],
    fs=11, bold0=False)

# ── Legend ────────────────────────────────────────────────────────────────────
lx, ly = 0.5, 9.0
ax.text(lx+0.1, ly+0.4, "Branch Color Key", fontsize=10,
        fontweight="bold", color=C_INPUT)
entries = [
    (C_CNN,  "CNN Branch  (EfficientNet-B4)"),
    (C_TEX,  "Texture Branch  (LBP + GLCM)"),
    (C_COL,  "Color Branch  (HSV Histogram)"),
    (C_FUSE, "AttentionFusion  (Phase 3)"),
    (C_HEAD, "Classification / Loss Heads"),
]
for i, (fc, lbl) in enumerate(entries):
    yy = ly - 0.55*(i+1)
    p = FancyBboxPatch((lx-0.02, yy-0.15), 0.38, 0.30,
                       boxstyle="round,pad=0", facecolor=fc,
                       edgecolor="white", linewidth=1, zorder=5)
    ax.add_patch(p)
    ax.text(lx+0.55, yy+0.015, lbl, fontsize=9, color=C_INPUT, va="center")

# ── Parameter Count ───────────────────────────────────────────────────────────
px, py = 0.5, 5.8
ax.text(px+0.1, py+0.4, "Parameter Count", fontsize=10,
        fontweight="bold", color=C_INPUT)
params = [
    ("EfficientNet-B4",    "~19.3 M", C_CNN),
    ("Texture + Color MLP","~56 K",   C_TEX),
    ("Fusion + All Heads", "~2.2 M",  C_FUSE),
    ("Total",              "~21.6 M", C_INPUT),
]
for i, (lbl, val, col) in enumerate(params):
    yy = py - 0.52*(i+1)
    fw = "bold" if lbl == "Total" else "normal"
    ax.text(px+0.1, yy,
            f"{lbl:<24s}  {val:>9s}",
            fontsize=8.5, color=col, fontweight=fw,
            fontfamily="monospace")

# ── Inference note ────────────────────────────────────────────────────────────
ax.text(0.6, 3.0, "Inference: ~2.70 ms/image (GPU)",
        fontsize=8.5, color=C_FUSE, fontstyle="italic",
        fontweight="bold")
ax.text(0.6, 2.55, "Dataset: SpiceSpectrum  |  11 classes  |  ~11,000 images",
        fontsize=8.0, color=C_INPUT)
ax.text(0.6, 2.1, "Split: 70% train / 10% val / 20% test  (seed=42)",
        fontsize=8.0, color=C_INPUT)

plt.tight_layout(pad=0.2)

out_dir = Path("D:/SpiceNet/outputs")
out_dir.mkdir(parents=True, exist_ok=True)
png_path  = out_dir / "architecture_diagram.png"
pptx_path = out_dir / "architecture_diagram.pptx"

plt.savefig(png_path, dpi=200, bbox_inches="tight",
            facecolor=C_BG, edgecolor="none")
plt.close()
print(f"[1/2] PNG  saved -> {png_path}")

# ── Build PPTX ────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(24)
prs.slide_height = Inches(15.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
fill  = slide.background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xF0, 0xF3, 0xF4)

slide.shapes.add_picture(
    str(png_path),
    left=Inches(0), top=Inches(0),
    width=Inches(24), height=Inches(15.5),
)

prs.save(str(pptx_path))
print(f"[2/2] PPTX saved -> {pptx_path}")
print("\nOpen architecture_diagram.pptx in PowerPoint.")
